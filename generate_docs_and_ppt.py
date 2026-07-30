import os
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml import parse_xml, OxmlElement
from docx.oxml.ns import nsdecls, qn

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR = r"c:\Users\User\pravallika sentinel"
IMG_DIR = os.path.join(BASE_DIR, "docs", "images")
DESKTOP_DIR = r"C:\Users\User\OneDrive\Desktop" if os.path.exists(r"C:\Users\User\OneDrive\Desktop") else r"C:\Users\User\Desktop"

# Image paths
IMG_DASHBOARD = os.path.join(IMG_DIR, "dashboard.png")
IMG_LOGIN = os.path.join(IMG_DIR, "login_page.png")
IMG_REGISTER = os.path.join(IMG_DIR, "registration_page.png")
IMG_CHAT = os.path.join(IMG_DIR, "ai_chat.png")
IMG_SCANNER = os.path.join(IMG_DIR, "file_scanner.png")

# ==============================================================================
# 1. BUILD WORD DOCUMENT (.DOCX)
# ==============================================================================
def create_word_document(save_paths):
    doc = docx.Document()
    
    # Page setup
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    # Styles
    normal_style = doc.styles['Normal']
    normal_style.font.name = 'Calibri'
    normal_style.font.size = Pt(11)
    normal_style.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Helper function for title & headers
    def add_title(text):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.name = 'Calibri'
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)
        return p

    def add_subtitle(text):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.name = 'Calibri'
        run.font.size = Pt(16)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        return p

    def add_heading1(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(18)
        p.paragraph_format.space_after = Pt(6)
        run = p.add_run(text)
        run.font.name = 'Calibri'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)
        return p

    def add_heading2(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after = Pt(4)
        run = p.add_run(text)
        run.font.name = 'Calibri'
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x66, 0x99)
        return p

    def add_paragraph(text, bold_prefix=None):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(6)
        p.paragraph_format.line_spacing = 1.15
        if bold_prefix:
            r_bold = p.add_run(bold_prefix)
            r_bold.bold = True
            r_bold.font.color.rgb = RGBColor(0x00, 0x33, 0x66)
        p.add_run(text)
        return p

    def add_image_with_caption(img_path, caption_text, width_inches=6.0):
        if os.path.exists(img_path):
            p_img = doc.add_paragraph()
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_img.paragraph_format.space_before = Pt(12)
            p_img.paragraph_format.space_after = Pt(4)
            p_img.add_run().add_picture(img_path, width=Inches(width_inches))

            p_cap = doc.add_paragraph()
            p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_cap.paragraph_format.space_after = Pt(14)
            r_cap = p_cap.add_run(f"Figure: {caption_text}")
            r_cap.font.italic = True
            r_cap.font.size = Pt(9.5)
            r_cap.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # --- COVER / TITLE ---
    add_title("PROJECT TECHNICAL REPORT")
    add_subtitle("SentinelGPT — Autonomous Cyber Defense Platform")
    doc.add_paragraph().paragraph_format.space_after = Pt(12)

    # Meta Table
    table = doc.add_table(rows=5, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    meta_data = [
        ("Project Title:", "SentinelGPT Autonomous Cyber Defense Platform"),
        ("Live Deployment Link:", "https://sentinelgpt-ai.vercel.app"),
        ("GitHub Repository:", "https://github.com/Pravallika2025/sentigraud-ai-.git"),
        ("Technology Stack:", "React 19, Vite 5.4, FastAPI, SQLAlchemy, SQLite, JWT"),
        ("Documentation Date:", "July 2026"),
    ]
    for idx, (k, v) in enumerate(meta_data):
        row = table.rows[idx]
        row.cells[0].paragraphs[0].add_run(k).bold = True
        row.cells[1].paragraphs[0].add_run(v)
    
    doc.add_page_break()

    # --- SECTION 1: EXECUTIVE SUMMARY & ABSTRACT ---
    add_heading1("1. Executive Summary & Abstract")
    add_paragraph("SentinelGPT is an enterprise-grade, full-stack autonomous Security Operations Center (SOC) platform designed for continuous network threat monitoring, real-time telemetry processing, heuristic vulnerability analysis, and automated quarantine response.")
    add_paragraph("Built to bridge real-time cybersecurity operations with cloud-native serverless architecture, SentinelGPT processes network anomalies, maps them dynamically against MITRE ATT&CK tactics, and renders actionable AI remediation recommendations for security analysts.")

    # --- SECTION 2: SYSTEM ARCHITECTURE & TECH STACK ---
    add_heading1("2. System Architecture & Technical Stack")
    add_paragraph("The application follows a decoupled monorepo architecture, leveraging React 19 for rendering dynamic glassmorphic telemetry components and FastAPI (Python) for asynchronous microservices and serverless database interaction.")
    
    add_heading2("Key Technology Components:")
    add_paragraph(" Modern single-page application framework providing fast Virtual DOM reconciliation and responsive glassmorphic dark theme components.", bold_prefix="• Frontend (React 19 + Vite 5.4): ")
    add_paragraph(" Asynchronous Python framework managing RESTful JSON APIs (`/api/*`), OpenAPI schemas, and WebSockets.", bold_prefix="• Backend (FastAPI + Python 3.11): ")
    add_paragraph(" SQLAlchemy ORM interfacing with SQLite (`/tmp/sentinel_vercel.db` on Vercel) for persistent incident logging.", bold_prefix="• Database & Storage: ")
    add_paragraph(" Expiring HS256 JWT tokens with SHA-256 password hashing and role-based access control.", bold_prefix="• Authentication: ")
    add_paragraph(" Multi-stage Vercel serverless deployment (`vercel.json`) serving static assets from `/dist` and dynamic API handlers.", bold_prefix="• Cloud Hosting: ")

    # --- SECTION 3: CORE OPERATIONS DASHBOARD ---
    add_heading1("3. Core Operations Dashboard Interface")
    add_paragraph("The main Security Operations Center (SOC) dashboard consolidates live perimeter telemetry into real-time KPI cards, threat velocity area charts, severity distribution donuts, sector heatmap matrices, and an automated quarantine management list.")
    
    add_image_with_caption(IMG_DASHBOARD, "SentinelGPT Real-time SOC Operations Dashboard Interface", width_inches=6.2)

    # --- SECTION 4: AUTHENTICATION & REGISTRATION ---
    add_heading1("4. Authentication & Operator Onboarding")
    add_paragraph("SentinelGPT incorporates a cyberpunk-themed authentication interface offering multi-role login, operator registration with organizational profile setup, and a 1-click Quick Demo Access bypass for rapid evaluation.")

    add_image_with_caption(IMG_LOGIN, "Cyberpunk Security Operator Login Interface with Quick Demo Access", width_inches=5.8)
    add_image_with_caption(IMG_REGISTER, "Operator Registration & Organization Setup Portal", width_inches=5.8)

    # --- SECTION 5: AI SECURITY ASSISTANT & FILE SCANNER ---
    add_heading1("5. Conversational AI Assistant & Payload Scanner")
    add_paragraph("To accelerate incident triage, SentinelGPT integrates a conversational AI security assistant trained on threat signatures and MITRE ATT&CK techniques, alongside a static payload scanner for analyzing suspicious log files and executables.")

    add_image_with_caption(IMG_CHAT, "Conversational AI Threat Triage Assistant Interface", width_inches=5.8)
    add_image_with_caption(IMG_SCANNER, "Heuristic Payload & Static Log File Scanner", width_inches=5.8)

    # --- SECTION 6: API SPECIFICATION & DEPLOYMENT ---
    add_heading1("6. API Endpoints & Live Deployment")
    add_paragraph("The backend exposes standard REST endpoints for telemetry queries, incident history export, and automated IP quarantine control:")
    
    api_table = doc.add_table(rows=6, cols=3)
    api_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    api_headers = ["Endpoint", "HTTP Method", "Functionality"]
    hdr_cells = api_table.rows[0].cells
    for i, title in enumerate(api_headers):
        hdr_cells[i].paragraphs[0].add_run(title).bold = True

    api_rows = [
        ("/api/login", "POST", "Authenticates operator, returns JWT token"),
        ("/api/snapshot", "GET", "Retrieves telemetry snapshot, logs & quarantine list"),
        ("/api/block_ip", "POST", "Adds target IP address to active firewall quarantine"),
        ("/api/unblock_ip", "POST", "Revokes quarantine for target IP address"),
        ("/health", "GET", "Returns live serverless system health status"),
    ]
    for r_idx, (ep, method, desc_text) in enumerate(api_rows, start=1):
        row_cells = api_table.rows[r_idx].cells
        row_cells[0].paragraphs[0].add_run(ep)
        row_cells[1].paragraphs[0].add_run(method)
        row_cells[2].paragraphs[0].add_run(desc_text)

    add_heading2("Live Deployment Details:")
    add_paragraph("https://sentinelgpt-ai.vercel.app", bold_prefix="• Production Web App: ")
    add_paragraph("https://sentinelgpt-ai.vercel.app/docs", bold_prefix="• Interactive Swagger API Docs: ")
    add_paragraph("https://github.com/Pravallika2025/sentigraud-ai-.git", bold_prefix="• GitHub Repository: ")

    # Save to all target paths
    for p in save_paths:
        doc.save(p)
        print(f"Word document saved to: {p}")

# ==============================================================================
# 2. BUILD POWERPOINT PRESENTATION (.PPTX)
# ==============================================================================
def create_ppt_presentation(save_paths):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)  # 16:9 Widescreen
    prs.slide_height = PptInches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Theme colors
    DARK_BG = PptRGBColor(0x08, 0x0C, 0x14)
    CARD_BG = PptRGBColor(0x0F, 0x17, 0x2A)
    CYAN_ACCENT = PptRGBColor(0x00, 0xF2, 0xFF)
    WHITE = PptRGBColor(0xFF, 0xFF, 0xFF)
    GRAY = PptRGBColor(0x94, 0xA3, 0xB8)
    DARK_BLUE = PptRGBColor(0x1E, 0x29, 0x3B)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, title_text, category_text="SENTINELGPT CYBER DEFENSE"):
        tb = slide.shapes.add_textbox(PptInches(0.8), PptInches(0.4), PptInches(11.7), PptInches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = PptPt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_ACCENT

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = PptPt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    # --- SLIDE 1: TITLE SLIDE ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    tb1 = slide1.shapes.add_textbox(PptInches(1.0), PptInches(2.0), PptInches(11.333), PptInches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "🛡️ SentinelGPT"
    p1.font.size = PptPt(48)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_ACCENT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "Autonomous Cyber Defense Platform"
    p2.font.size = PptPt(28)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf1.add_paragraph()
    p3.text = "\nReal-time AI Security Monitoring • Heuristic Anomaly Analysis • Autonomous Quarantine"
    p3.font.size = PptPt(16)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf1.add_paragraph()
    p4.text = "\n🌐 Live Demo: https://sentinelgpt-ai.vercel.app  |  📁 GitHub: Pravallika2025/sentigraud-ai-"
    p4.font.size = PptPt(13)
    p4.font.color.rgb = CYAN_ACCENT
    p4.alignment = PP_ALIGN.CENTER

    # --- SLIDE 2: PROBLEM & SOLUTION ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Problem Statement & Proposed Solution")

    tb2_left = slide2.shapes.add_textbox(PptInches(0.8), PptInches(1.6), PptInches(5.6), PptInches(5.2))
    tf2_l = tb2_left.text_frame
    tf2_l.word_wrap = True
    tf2_l.paragraphs[0].text = "🔴 Traditional SOC Challenges:"
    tf2_l.paragraphs[0].font.size = PptPt(18)
    tf2_l.paragraphs[0].font.bold = True
    tf2_l.paragraphs[0].font.color.rgb = PptRGBColor(255, 100, 100)

    challenges = [
        "High alert volume causing security analyst fatigue.",
        "Delayed manual intervention leading to breach escalation.",
        "Fragmented tools lacking unified MITRE ATT&CK correlation.",
        "Lack of real-time automated quarantine response.",
    ]
    for c in challenges:
        p = tf2_l.add_paragraph()
        p.text = f"• {c}"
        p.font.size = PptPt(14)
        p.font.color.rgb = GRAY

    tb2_right = slide2.shapes.add_textbox(PptInches(6.8), PptInches(1.6), PptInches(5.7), PptInches(5.2))
    tf2_r = tb2_right.text_frame
    tf2_r.word_wrap = True
    tf2_r.paragraphs[0].text = "🟢 SentinelGPT Solution:"
    tf2_r.paragraphs[0].font.size = PptPt(18)
    tf2_r.paragraphs[0].font.bold = True
    tf2_r.paragraphs[0].font.color.rgb = CYAN_ACCENT

    solutions = [
        "Autonomous heuristic engine running 24/7 background telemetry.",
        "Instant visual triage via dynamic glassmorphism dashboard.",
        "Automated IP quarantine with 1-click revoke controls.",
        "AI assistant generating instant mitigation recommendations.",
    ]
    for s in solutions:
        p = tf2_r.add_paragraph()
        p.text = f"• {s}"
        p.font.size = PptPt(14)
        p.font.color.rgb = WHITE

    # --- SLIDE 3: SYSTEM ARCHITECTURE ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "System Architecture & Technology Stack")

    tb3 = slide3.shapes.add_textbox(PptInches(0.8), PptInches(1.6), PptInches(11.7), PptInches(5.2))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    tech_stack = [
        ("React 19 + Vite 5.4 SPA", "Delivers ultra-responsive glassmorphism SOC interface with dark theme neon accents."),
        ("FastAPI Python Backend", "Asynchronous microservices handling REST API, OpenAPI docs, and WebSockets."),
        ("SQLAlchemy & SQLite DB", "Persistent threat logs, blocked IPs quarantine list, and user profiles."),
        ("JWT Authentication", "Secure 8-hour expiring HS256 tokens with SHA-256 hashed user credentials."),
        ("Vercel Serverless Platform", "Full-stack serverless cloud deployment serving static SPA build + Python functions."),
    ]
    for idx, (tech, desc_text) in enumerate(tech_stack):
        p = tf3.add_paragraph() if idx > 0 else tf3.paragraphs[0]
        p.text = f"🔹 {tech}: "
        p.font.size = PptPt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
        
        p_desc = tf3.add_paragraph()
        p_desc.text = f"    {desc_text}\n"
        p_desc.font.size = PptPt(13)
        p_desc.font.color.rgb = GRAY

    # --- SLIDE 4: DASHBOARD SHOWCASE ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Real-Time SOC Operations Dashboard")

    if os.path.exists(IMG_DASHBOARD):
        slide4.shapes.add_picture(IMG_DASHBOARD, PptInches(0.8), PptInches(1.6), width=PptInches(11.7))

    # --- SLIDE 5: AUTH & ONBOARDING ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Authentication & Operator Onboarding")

    if os.path.exists(IMG_LOGIN):
        slide5.shapes.add_picture(IMG_LOGIN, PptInches(0.8), PptInches(1.6), width=PptInches(5.6))
    if os.path.exists(IMG_REGISTER):
        slide5.shapes.add_picture(IMG_REGISTER, PptInches(6.8), PptInches(1.6), width=PptInches(5.6))

    # --- SLIDE 6: AI ASSISTANT & SCANNER ---
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "AI Threat Triage & Payload File Scanner")

    if os.path.exists(IMG_CHAT):
        slide6.shapes.add_picture(IMG_CHAT, PptInches(0.8), PptInches(1.6), width=PptInches(5.6))
    if os.path.exists(IMG_SCANNER):
        slide6.shapes.add_picture(IMG_SCANNER, PptInches(6.8), PptInches(1.6), width=PptInches(5.6))

    # --- SLIDE 7: CONCLUSION & ACCESS ---
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Project Submission & Access Links")

    tb7 = slide7.shapes.add_textbox(PptInches(1.5), PptInches(2.0), PptInches(10.333), PptInches(4.5))
    tf7 = tb7.text_frame
    tf7.word_wrap = True

    links = [
        ("🌐 Official Live Web Deployment:", "https://sentinelgpt-ai.vercel.app"),
        ("📡 Live Server Health Status:", "https://sentinelgpt-ai.vercel.app/health"),
        ("📖 Interactive Swagger API Docs:", "https://sentinelgpt-ai.vercel.app/docs"),
        ("📁 Official GitHub Repository:", "https://github.com/Pravallika2025/sentigraud-ai-.git"),
    ]
    for label, url in links:
        p1_l = tf7.add_paragraph()
        p1_l.text = label
        p1_l.font.size = PptPt(16)
        p1_l.font.bold = True
        p1_l.font.color.rgb = WHITE

        p2_u = tf7.add_paragraph()
        p2_u.text = f"👉 {url}\n"
        p2_u.font.size = PptPt(15)
        p2_u.font.color.rgb = CYAN_ACCENT

    # Save PPT
    for p in save_paths:
        prs.save(p)
        print(f"PowerPoint saved to: {p}")

if __name__ == "__main__":
    doc_paths = [
        os.path.join(BASE_DIR, "docs", "SentinelGPT_Project_Report.docx"),
        os.path.join(DESKTOP_DIR, "SentinelGPT_Project_Report.docx")
    ]
    ppt_paths = [
        os.path.join(BASE_DIR, "docs", "SentinelGPT_Project_Presentation.pptx"),
        os.path.join(DESKTOP_DIR, "SentinelGPT_Project_Presentation.pptx")
    ]

    create_word_document(doc_paths)
    create_ppt_presentation(ppt_paths)
